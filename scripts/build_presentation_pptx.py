"""기말과제 발표자료 .pptx — 베이지 클래식 디자인 (수정 2).

사용자 주문 (2차):
1. 각 페이지 제목을 제외한 텍스트 30% 키움
2. 페이지 제목 밑 구분선을 제목 바로 밑으로 올림 (1.75 → 1.6)
3. 모든 박스를 둥근 모서리로 교체 (RECTANGLE → ROUNDED_RECTANGLE)
4. 마지막 페이지 (NEXT) 레이아웃 재배치 — LIMITATION 도 박스 안으로
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


OUTPUT = Path("이용주-기말발표자료.pptx")

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── 베이지 클래식 팔레트 ─────────────────────────────────
BG_BEIGE = RGBColor(0xF5, 0xF0, 0xE8)
CARD_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BORDER = RGBColor(0xE5, 0xDD, 0xD0)
SHADOW = RGBColor(0xD5, 0xC9, 0xB4)
DIVIDER = RGBColor(0xD5, 0xCC, 0xB8)

TEXT_DARK = RGBColor(0x2C, 0x24, 0x18)
TEXT_GRAY = RGBColor(0x6B, 0x63, 0x56)
TEXT_FAINT = RGBColor(0x9B, 0x91, 0x82)

ACCENT_GREEN = RGBColor(0x2D, 0x4A, 0x3E)
ACCENT_WINE = RGBColor(0x8B, 0x26, 0x35)
ACCENT_GOLD = RGBColor(0xB5, 0x79, 0x3E)

INFO_BG = RGBColor(0xF9, 0xF5, 0xEE)

CORNER = 0.06  # 둥근 정도 (모든 카드 통일)


def set_run_font(run, size=14, bold=False, color=None, name="맑은 고딕"):
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", name)
    run.font.size = Pt(size)
    run.bold = bold
    if color:
        run.font.color.rgb = color


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=None,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=None):
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    items = text if isinstance(text, list) else [text]
    for i, line in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        set_run_font(r, size=size, bold=bold, color=color)
    return tx


def add_rect(slide, x, y, w, h, fill=None, line=None, line_width=0.75,
              rounded=True):
    """기본적으로 둥근 사각형. divider 같이 얇은 라인엔 rounded=False."""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.shadow.inherit = False
    if rounded:
        try:
            shape.adjustments[0] = CORNER
        except Exception:
            pass
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width)
    return shape


def add_slide_bg(slide):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=BG_BEIGE, rounded=False)


def add_card(slide, x, y, w, h, fill=None, with_shadow=True, corner=None):
    """그림자 + 둥근 흰 카드."""
    fill = fill or CARD_WHITE
    corner = corner if corner is not None else CORNER
    if with_shadow:
        offset = 0.04
        sh = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x + offset), Inches(y + offset),
            Inches(w), Inches(h)
        )
        try:
            sh.adjustments[0] = corner
        except Exception:
            pass
        sh.fill.solid()
        sh.fill.fore_color.rgb = SHADOW
        sh.line.fill.background()
        sh.shadow.inherit = False
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    try:
        card.adjustments[0] = corner
    except Exception:
        pass
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    card.line.color.rgb = CARD_BORDER
    card.line.width = Pt(0.75)
    card.shadow.inherit = False
    return card


def add_slide_header(slide, eng, kor, slide_num, kor_size=22):
    """헤더 — 제목 + 제목 바로 밑 구분선."""
    add_text(slide, Inches(0.7), Inches(0.55), Inches(8), Inches(0.35),
             eng.upper(), size=11, bold=True, color=ACCENT_GREEN)
    add_text(slide, Inches(0.7), Inches(0.92), Inches(11.9), Inches(0.7),
             kor, size=kor_size, bold=True, color=TEXT_DARK)
    # 구분선 — 제목 바로 밑 (1.75 → 1.62)
    add_rect(slide, Inches(0.7), Inches(1.62),
             Inches(11.9), Inches(0.015), fill=DIVIDER, rounded=False)
    # 우하단 슬라이드 번호 (모노 스타일)
    add_text(slide, Inches(12.2), Inches(7.05), Inches(0.7), Inches(0.3),
             f"{slide_num:02d} / 09", size=11.5, color=TEXT_FAINT,
             align=PP_ALIGN.RIGHT)


def add_kicker(slide, x, y, w, h, label, color=None):
    """작은 영문 라벨 — 폰트 크기 30% 키움 (10 → 13)."""
    color = color or ACCENT_GREEN
    add_text(slide, Inches(x), Inches(y), Inches(w), Inches(h),
             label.upper(), size=13, bold=True, color=color)


def add_icon_card(slide, x, y, w, h, title, body, icon_color=None,
                  title_size=16, body_size=14):
    """둥근 카드 + 좌상단 작은 정사각 아이콘 + 제목 + 본문 (30% 키움)."""
    icon_color = icon_color or ACCENT_GREEN
    add_card(slide, x, y, w, h)
    # 아이콘 (작은 그린 둥근 정사각형)
    icon = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x + 0.3), Inches(y + 0.3),
        Inches(0.4), Inches(0.4)
    )
    try:
        icon.adjustments[0] = 0.20
    except Exception:
        pass
    icon.fill.solid()
    icon.fill.fore_color.rgb = icon_color
    icon.line.fill.background()
    icon.shadow.inherit = False
    # 제목
    add_text(slide, Inches(x + 0.85), Inches(y + 0.3), Inches(w - 1.1),
             Inches(0.45), title, size=title_size, bold=True, color=TEXT_DARK,
             anchor=MSO_ANCHOR.MIDDLE)
    # 본문 (30% 키움)
    body_items = body if isinstance(body, list) else [body]
    add_text(slide, Inches(x + 0.3), Inches(y + 0.95), Inches(w - 0.6),
             Inches(h - 1.1), body_items, size=body_size, color=TEXT_GRAY,
             line_spacing=1.4)


def add_row_card(slide, x, y, w, h, label, content,
                 label_size=15, content_size=14):
    """표 행 카드 — 라벨/내용 모두 30% 키움."""
    add_card(slide, x, y, w, h)
    # 좌측 작은 그린 마커 (둥근 정사각형)
    marker = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x + 0.3), Inches(y + h / 2 - 0.06),
        Inches(0.12), Inches(0.12)
    )
    try:
        marker.adjustments[0] = 0.35
    except Exception:
        pass
    marker.fill.solid()
    marker.fill.fore_color.rgb = ACCENT_GREEN
    marker.line.fill.background()
    marker.shadow.inherit = False
    # 라벨
    add_text(slide, Inches(x + 0.6), Inches(y), Inches(3.0), Inches(h),
             label, size=label_size, bold=True, color=TEXT_DARK,
             anchor=MSO_ANCHOR.MIDDLE)
    # 내용
    add_text(slide, Inches(x + 3.8), Inches(y), Inches(w - 4.0), Inches(h),
             content, size=content_size, color=TEXT_GRAY,
             anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)


def add_flow_step(slide, x, y, w, h, title, sub=""):
    """둥근 카드 단계 — 폰트 키움."""
    add_card(slide, x, y, w, h)
    add_text(slide, Inches(x), Inches(y + 0.2), Inches(w), Inches(0.4),
             title, size=15.5, bold=True, color=ACCENT_GREEN,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    if sub:
        add_text(slide, Inches(x), Inches(y + 0.6), Inches(w), Inches(h - 0.65),
                 sub, size=12, color=TEXT_GRAY, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.TOP, line_spacing=1.2)


def add_arrow(slide, x, y):
    arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                  Inches(x), Inches(y), Inches(0.3), Inches(0.36))
    arr.fill.solid()
    arr.fill.fore_color.rgb = ACCENT_GREEN
    arr.line.fill.background()
    arr.shadow.inherit = False


def add_hyperlink_text(slide, x, y, w, h, text, url, size=14,
                        color=None, bold=False, align=PP_ALIGN.LEFT):
    color = color or ACCENT_GOLD
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    set_run_font(r, size=size, bold=bold, color=color)
    r.font.underline = True
    r.hyperlink.address = url
    return tx


def make_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ────────────────────────────────────────────
    # Slide 1 — 표지
    # 폰트 키움: 한 줄 소개 17→22, 발표자/링크 11→14, 라벨 8.5→11
    # ────────────────────────────────────────────
    s = make_slide(prs)
    add_slide_bg(s)
    # 우측 장식 (둥근 형태 유지 — 원)
    deco1 = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                Inches(10.3), Inches(-1.5),
                                Inches(5.0), Inches(5.0))
    deco1.fill.solid()
    deco1.fill.fore_color.rgb = ACCENT_GREEN
    deco1.line.fill.background()
    deco1.shadow.inherit = False
    deco2 = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                Inches(11.5), Inches(4.2),
                                Inches(3.2), Inches(3.2))
    deco2.fill.solid()
    deco2.fill.fore_color.rgb = ACCENT_GOLD
    deco2.line.fill.background()
    deco2.shadow.inherit = False

    # 좌상단 강의 칩 (둥근)
    chip = slide_chip = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.7), Inches(0.7), Inches(3.6), Inches(0.45)
    )
    try:
        chip.adjustments[0] = 0.45
    except Exception:
        pass
    chip.fill.solid()
    chip.fill.fore_color.rgb = CARD_WHITE
    chip.line.color.rgb = CARD_BORDER
    chip.line.width = Pt(0.75)
    chip.shadow.inherit = False
    add_text(s, Inches(0.85), Inches(0.7), Inches(3.4), Inches(0.45),
             "BAVECODING  ·  바이브코딩 기말과제", size=12, bold=True,
             color=ACCENT_GREEN, anchor=MSO_ANCHOR.MIDDLE)

    # 브랜드 마크 (둥근 그린 박스 + T)
    mark = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.7), Inches(2.4),
                               Inches(0.95), Inches(0.95))
    try:
        mark.adjustments[0] = 0.18
    except Exception:
        pass
    mark.fill.solid()
    mark.fill.fore_color.rgb = ACCENT_GREEN
    mark.line.fill.background()
    mark.shadow.inherit = False
    add_text(s, Inches(0.7), Inches(2.4), Inches(0.95), Inches(0.95),
             "T", size=32, bold=True, color=CARD_WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.8), Inches(2.55), Inches(6), Inches(0.5),
             "CONTENT  DIGEST", size=14, bold=True, color=TEXT_FAINT)

    # 큰 제목 (그대로 — 페이지 제목)
    add_text(s, Inches(0.7), Inches(3.55), Inches(11), Inches(1.0),
             "TubeBrief", size=42, bold=True, color=TEXT_DARK)
    # 부제 (그대로 — 제목의 일부)
    add_text(s, Inches(0.7), Inches(4.55), Inches(11), Inches(0.5),
             "유튜브 구독 채널 자동 요약 서비스",
             size=18, bold=True, color=ACCENT_GREEN)

    # 한 줄 소개 (17→22, 30% 키움)
    add_text(s, Inches(0.7), Inches(5.25), Inches(9.0), Inches(1.2),
             "구독한 유튜브 채널의 신규 영상을 매일 자동으로 감지해 AI로 구조화된 "
             "한 장짜리 요약 보고서를 작성·아카이빙하는 개인용 콘텐츠 다이제스트.",
             size=22, color=TEXT_GRAY, line_spacing=1.5)

    # 구분선
    add_rect(s, Inches(0.7), Inches(6.55), Inches(11.9), Inches(0.015),
             fill=DIVIDER, rounded=False)

    # 발표자 (라벨 8.5→11, 정보 11→14)
    add_text(s, Inches(0.7), Inches(6.75), Inches(2), Inches(0.25),
             "PRESENTER", size=11, bold=True, color=TEXT_FAINT)
    add_text(s, Inches(0.7), Inches(7.0), Inches(4), Inches(0.3),
             "이용주  ·  학번: _(제출 전 본인 학번 기입)_",
             size=14, color=TEXT_DARK)

    # GitHub
    add_text(s, Inches(4.5), Inches(6.75), Inches(2), Inches(0.25),
             "GITHUB", size=11, bold=True, color=TEXT_FAINT)
    add_hyperlink_text(s, Inches(4.5), Inches(7.0), Inches(5), Inches(0.3),
                       "github.com/leeyongjoo9-rgb/tubebrief_mvp",
                       "https://github.com/leeyongjoo9-rgb/tubebrief_mvp",
                       size=14)

    # Service URL
    add_text(s, Inches(9.0), Inches(6.75), Inches(2), Inches(0.25),
             "SERVICE", size=11, bold=True, color=TEXT_FAINT)
    add_hyperlink_text(s, Inches(9.0), Inches(7.0), Inches(4), Inches(0.3),
                       "tubebrief-mvp.vercel.app",
                       "https://tubebrief-mvp.vercel.app/",
                       size=14)

    # ────────────────────────────────────────────
    # Slide 2 — WHY
    # ────────────────────────────────────────────
    s = make_slide(prs)
    add_slide_bg(s)
    add_slide_header(s, "WHY  ·  문제정의와 대상 사용자",
                     "정보는 쏟아지고 시간은 부족하다.", 2, kor_size=22)

    # 좌측 PROBLEM 카드
    add_card(s, x=0.7, y=2.0, w=5.95, h=5.0)
    add_kicker(s, x=0.95, y=2.25, w=4, h=0.3,
               label="PROBLEM  ·  문제 · 사용 상황", color=ACCENT_WINE)
    add_text(s, Inches(0.95), Inches(2.65), Inches(5.5), Inches(0.45),
             "이런 불편을 매일 겪는다", size=19.5, bold=True, color=TEXT_DARK)
    add_text(s, Inches(0.95), Inches(3.4), Inches(5.5), Inches(3.5),
             ["•  매일 채널마다 1~5편 신규 영상 누적, 다 볼 시간 없음",
              "•  어떤 영상에 시간을 투자할지 우선순위 판단 어려움",
              "•  유튜브 자막은 길고 비구조화 텍스트라 훑어 읽기 어려움",
              "•  한 줄 캡션은 정보 밀도 부족, 알고리즘 추천은 관심사와 어긋남",
              "•  오래된 영상은 RSS 에서 사라져 회수할 방법 없음"],
             size=15.5, color=TEXT_GRAY, line_spacing=1.6)

    # 우측 GOAL 카드
    add_card(s, x=6.85, y=2.0, w=5.95, h=5.0)
    add_kicker(s, x=7.1, y=2.25, w=4, h=0.3,
               label="GOAL  ·  대상 사용자 · 해결 가설")
    add_text(s, Inches(7.1), Inches(2.65), Inches(5.5), Inches(0.45),
             "이런 사용자를 위해", size=19.5, bold=True, color=TEXT_DARK)
    add_text(s, Inches(7.1), Inches(3.4), Inches(5.5), Inches(2.2),
             ["•  채널을 등록만 해두면 매일 자동으로 요약 보고서 누적",
              "•  영상 한 편을 1분 안에 파악 → 골라서 시청",
              "•  단순 캡션 X → 헤드라인·본문·출연자·기업·시각 점프까지 구조화"],
             size=15.5, color=TEXT_GRAY, line_spacing=1.6)
    # 강조 박스 (둥근)
    add_rect(s, Inches(7.1), Inches(5.65), Inches(5.55), Inches(1.3),
             fill=INFO_BG, line=CARD_BORDER, line_width=0.5)
    add_text(s, Inches(7.3), Inches(5.8), Inches(5.2), Inches(1.0),
             ["▸  구독 채널 3개 이상, 콘텐츠 중심 시청 / 시간 부족한 개인",
              "▸  정보 비용 최소화로 의사결정에 쓰고 싶은 사용자"],
             size=15, color=TEXT_DARK, line_spacing=1.55, bold=True)

    # ────────────────────────────────────────────
    # Slide 3 — WHAT
    # ────────────────────────────────────────────
    s = make_slide(prs)
    add_slide_bg(s)
    add_slide_header(s, "WHAT  ·  주요 기능",
                     "수집부터 요약·구독·자동화까지", 3, kor_size=22)
    cards = [
        ("자동 수집", "채널 RSS 주기 폴링으로 신규 영상 자동 감지 (UNIQUE 로 중복 방지)"),
        ("자막 추출", "한국어 자막 추출, 실패 시 RSS 설명문 메타데이터로 폴백"),
        ("AI 구조화 요약",
         "gpt-5-mini structured outputs — 헤드라인·본문·주제·인물·기업·타임스탬프 분리"),
        ("유연한 채널 등록",
         "채널 전체 또는 제목·설명문 키워드 AND 필터로 코너 단위 수신"),
        ("출연자 ↔ 거론 인물 분리",
         "people / mentioned_people 별도 필드. 호스트 자동 제외, 환각 방지"),
        ("자동화 + 보고서 정리",
         "Vercel Cron 일 1회. 한 번 본 보고서는 소프트 삭제로 정리"),
    ]
    for i, (title, body) in enumerate(cards):
        row = i // 3
        col = i % 3
        x = 0.7 + col * 4.10
        y = 1.95 + row * 2.50
        add_icon_card(s, x=x, y=y, w=3.90, h=2.25, title=title, body=body,
                      title_size=16, body_size=14)

    # ────────────────────────────────────────────
    # Slide 4 — USE
    # ────────────────────────────────────────────
    s = make_slide(prs)
    add_slide_bg(s)
    add_slide_header(s, "USE  ·  사용자 이용 흐름",
                     "채널 등록 → 자동 처리 → 결과 조회", 4, kor_size=22)

    # INPUT
    add_kicker(s, x=0.7, y=2.05, w=4, h=0.3,
               label="INPUT  ·  입력 (1회 등록)", color=ACCENT_WINE)
    add_text(s, Inches(0.7), Inches(2.45), Inches(11.9), Inches(0.6),
             "사용자는 메인화면 [채널 관리] 에서 채널 URL + 키워드 필터를 1회 등록한다.",
             size=17, color=TEXT_DARK, line_spacing=1.4)

    # AUTO
    add_kicker(s, x=0.7, y=3.4, w=5, h=0.3,
               label="AUTO  ·  자동 처리 (매일 0시 UTC)")
    steps = [
        ("매일 0시 UTC", "Vercel Cron"),
        ("RSS 폴링", "구독·키워드 매치"),
        ("자막 추출", "youtube-transcript\n→ 폴백"),
        ("LLM 요약", "gpt-5-mini\nJSON Schema"),
        ("Supabase 저장", "videos.summary\njsonb"),
    ]
    cur_x = 0.7
    step_w = 2.30
    gap = 0.12
    for i, (t, sub) in enumerate(steps):
        add_flow_step(s, x=cur_x, y=3.85, w=step_w, h=1.25, title=t, sub=sub)
        if i < len(steps) - 1:
            add_arrow(s, cur_x + step_w + 0.04, 4.35)
        cur_x += step_w + gap

    # OUTPUT
    add_kicker(s, x=0.7, y=5.45, w=3, h=0.3, label="OUTPUT  ·  결과")
    add_card(s, x=0.7, y=5.85, w=5.9, h=1.4)
    add_text(s, Inches(0.95), Inches(6.0), Inches(5.4), Inches(1.15),
             ["▸  홈 — ChannelStrip, 카드 그리드, 마지막 확인 시각",
              "▸  채널별 필터링, 최신 영상이 좌측 첫 카드부터"],
             size=15, color=TEXT_DARK, line_spacing=1.45)
    add_card(s, x=6.85, y=5.85, w=5.9, h=1.4)
    add_text(s, Inches(7.1), Inches(6.0), Inches(5.4), Inches(1.15),
             ["▸  상세 — 헤드라인 + 키워드 + 출연자 / 거론 인물",
              "▸  (mm:ss) → 유튜브 그 시점 점프. 보고서 삭제 가능"],
             size=15, color=TEXT_DARK, line_spacing=1.45)

    # ────────────────────────────────────────────
    # Slide 5 — DEMO (큰 박스 1개)
    # ────────────────────────────────────────────
    s = make_slide(prs)
    add_slide_bg(s)
    add_slide_header(s, "DEMO  ·  데모 화면 (라이브)",
                     "tubebrief-mvp.vercel.app", 5, kor_size=19)

    add_kicker(s, x=0.7, y=1.95, w=3, h=0.3, label="LIVE", color=ACCENT_GOLD)
    add_hyperlink_text(s, Inches(1.3), Inches(1.95), Inches(8), Inches(0.35),
                       "https://tubebrief-mvp.vercel.app/",
                       "https://tubebrief-mvp.vercel.app/",
                       size=15, bold=True)
    add_text(s, Inches(0.7), Inches(2.4), Inches(11.9), Inches(0.4),
             "URL 클릭 → 라이브 사이트. 접속 불가 시 다음 슬라이드의 캡처 본 페이지로 이동.",
             size=13, color=TEXT_FAINT)

    # 큰 메인 캡처 카드 1개
    add_card(s, x=0.7, y=2.95, w=11.9, h=4.0)
    add_text(s, Inches(0.7), Inches(4.6), Inches(11.9), Inches(0.5),
             "[ 메인 화면 캡처 — 홈 대시보드 ]",
             size=18, bold=True, color=TEXT_FAINT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(0.7), Inches(5.1), Inches(11.9), Inches(0.4),
             "헤더 + ChannelStrip + 카드 그리드 (55개 요약 영상)",
             size=14, color=TEXT_FAINT,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ────────────────────────────────────────────
    # Slide 6 — STACK
    # ────────────────────────────────────────────
    s = make_slide(prs)
    add_slide_bg(s)
    add_slide_header(s, "STACK  ·  기술 스택 및 구현 구조",
                     "Next.js 풀스택 + Supabase + OpenAI", 6, kor_size=22)

    add_kicker(s, x=0.7, y=2.0, w=5, h=0.3,
               label="PIPELINE  ·  4단계 직렬 처리")
    pipeline = [("RSS 감지", ""), ("자막/폴백", ""), ("AI 요약", ""), ("DB 저장", "")]
    cur_x = 0.7
    step_w = 2.80
    for i, (t, sub) in enumerate(pipeline):
        add_flow_step(s, x=cur_x, y=2.4, w=step_w, h=1.1, title=t, sub=sub)
        if i < len(pipeline) - 1:
            add_arrow(s, cur_x + step_w + 0.04, 2.78)
        cur_x += step_w + 0.15

    rows = [
        ("프레임워크 / 언어",
         "Next.js 16 (App Router) · TypeScript · Tailwind v4 + shadcn/ui"),
        ("데이터 / 모델",
         "Supabase (PostgreSQL · jsonb) · OpenAI gpt-5-mini (JSON 스키마 strict)"),
        ("외부 데이터 소스",
         "YouTube RSS · 영상 페이지 스크래핑 · youtube-transcript 라이브러리"),
        ("배포 / 운영",
         "Vercel (Hobby) · Vercel Cron 일 1회 · GitHub 자동 빌드 연동"),
    ]
    for i, (label, content) in enumerate(rows):
        y = 4.0 + i * 0.80
        add_row_card(s, x=0.7, y=y, w=11.9, h=0.70,
                     label=label, content=content,
                     label_size=15, content_size=14)

    # ────────────────────────────────────────────
    # Slide 7 — API
    # ────────────────────────────────────────────
    s = make_slide(prs)
    add_slide_bg(s)
    add_slide_header(s, "API  ·  외부 API · MCP 활용 방식",
                     "직접 호출 4 + 인프라 활용 1", 7, kor_size=22)
    rows = [
        ("OpenAI Chat (gpt-5-mini)",
         "자막 + 시스템 프롬프트 + JSON Schema → strict mode → Supabase jsonb"),
        ("Supabase REST + JS SDK",
         "서버는 service role, 클라이언트는 publishable key (RLS 보호)"),
        ("YouTube RSS",
         "feeds/videos.xml?channel_id 폴링 → XML 파싱 → 키워드 매치 → INSERT"),
        ("YouTube 페이지 스크래핑",
         "URL → channel_id 추출. 6 패턴 fallback (canonical/og:url/Schema.org/JSON)"),
        ("Vercel Platform (인프라)",
         "vercel.json crons → 매일 0시 /api/cron. push 시 자동 재배포"),
    ]
    for i, (label, content) in enumerate(rows):
        y = 2.05 + i * 0.92
        add_row_card(s, x=0.7, y=y, w=11.9, h=0.80,
                     label=label, content=content,
                     label_size=15, content_size=14)
    add_text(s, Inches(0.7), Inches(6.9), Inches(11.9), Inches(0.35),
             "한계  ·  LLM 비용 영상당 ~$0.025 / Vercel Hobby 60초 timeout "
             "→ 보수적 limit + temperature 0.2 로 retry 최소화",
             size=12.5, color=TEXT_FAINT)

    # ────────────────────────────────────────────
    # Slide 8 — CONTEXT
    # ────────────────────────────────────────────
    s = make_slide(prs)
    add_slide_bg(s)
    add_slide_header(s, "CONTEXT  ·  컨텍스트 엔지니어링 활용",
                     ".md 파일 4개로 역할 · 목표 · 제약 · 검증 기준 설계", 8,
                     kor_size=19)
    rows = [
        ("CLAUDE.md",
         "세션마다 자동 로드. 절대 규칙 6개 + 작업 순서 8단계 + 코딩 스타일"),
        ("TubeBrief_PRD.md",
         "문제·기능·사용자 흐름·파이프라인. 자막 폴백 3단계가 그대로 status 값"),
        ("SUMMARY_SCHEMA.json",
         "LLM 출력 JSON 스키마 진실 공급원. OpenAI response_format 에 그대로 전달"),
        ("PRESENTATION.md",
         "발표 콘텐츠 + 옛/새 규칙 비교 라이브 자료. GitHub 리뷰 시 노출"),
    ]
    for i, (label, content) in enumerate(rows):
        y = 2.05 + i * 0.92
        add_row_card(s, x=0.7, y=y, w=11.9, h=0.80,
                     label=label, content=content,
                     label_size=15, content_size=14)
    add_kicker(s, x=0.7, y=5.85, w=5, h=0.3,
               label="INPUT  ·  Context 작성을 위한 자료")
    add_text(s, Inches(0.7), Inches(6.25), Inches(11.9), Inches(1.1),
             ["•  강의 교안의 'AI 협업 규칙' + Anthropic Claude Code 공식 베스트 프랙티스 → CLAUDE.md",
              "•  본인 문제 정의 + 기존 유튜브 요약 서비스 한계 분석 → PRD",
              "•  OpenAI structured outputs 공식 가이드 + JSON Schema draft-07 → SUMMARY_SCHEMA"],
             size=13, color=TEXT_GRAY, line_spacing=1.5)

    # ────────────────────────────────────────────
    # Slide 9 — NEXT (재배치: LIMITATION 도 박스 안)
    # ────────────────────────────────────────────
    s = make_slide(prs)
    add_slide_bg(s)
    add_slide_header(s, "NEXT  ·  어려움 · 해결 · 한계 · 개선",
                     "3겹 가드레일이 LLM 시스템의 핵심이었다.", 9, kor_size=19)

    # ── 상단: 좌측 어려움+해결 (높이 줄임, 3.4 → 3.4 유지하되 위치)
    # 좌측 박스
    add_card(s, x=0.7, y=2.0, w=6.1, h=3.4)
    add_kicker(s, x=0.95, y=2.2, w=4, h=0.3,
               label="CHALLENGE  ·  어려움", color=ACCENT_WINE)
    add_text(s, Inches(0.95), Inches(2.55), Inches(5.8), Inches(0.45),
             "LLM 한국어 요약이 규칙을 안 지킴", size=17.5, bold=True,
             color=TEXT_DARK)
    add_text(s, Inches(0.95), Inches(3.05), Inches(5.8), Inches(0.8),
             ["•  초기 gpt-4o-mini: brief 290자 / 1문단 / (mm:ss) 0개",
              "•  프롬프트 '가이드' → LLM 이 권고로 받아들임"],
             size=14, color=TEXT_GRAY, line_spacing=1.4)
    add_kicker(s, x=0.95, y=3.95, w=4, h=0.3,
               label="SOLUTION  ·  3겹 가드레일")
    add_text(s, Inches(0.95), Inches(4.3), Inches(5.8), Inches(1.05),
             ["①  프롬프트 명령조 강화 (반드시 N자 이상)",
              "②  validateSummary 후처리 + 위반 시 자동 재시도 (MAX 2)",
              "③  gpt-5-mini 로 모델 교체"],
             size=14.5, color=TEXT_DARK, line_spacing=1.45, bold=True)

    # 우측 결과 비교 박스
    add_card(s, x=7.0, y=2.0, w=5.7, h=3.4)
    add_kicker(s, x=7.25, y=2.2, w=4, h=0.3,
               label="RESULT  ·  같은 영상, 8배 정보 밀도")
    # 표 헤더
    th_y = 2.65
    add_rect(s, Inches(7.25), Inches(th_y), Inches(5.2), Inches(0.4),
             fill=INFO_BG)
    add_text(s, Inches(7.4), Inches(th_y), Inches(1.9), Inches(0.4),
             "항목", size=13.5, bold=True, color=TEXT_DARK,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(9.3), Inches(th_y), Inches(1.5), Inches(0.4),
             "옛 규칙", size=13.5, bold=True, color=TEXT_GRAY,
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(10.8), Inches(th_y), Inches(1.6), Inches(0.4),
             "새 규칙", size=13.5, bold=True, color=ACCENT_GREEN,
             anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    table_rows = [
        ("brief 분량",       "264자 / 1문단",   "2,122자 / 10문단"),
        ("(mm:ss) 마커",     "0",              "5"),
        ("companies",        "0",              "7"),
        ("출연자 / 거론",     "혼재",           "people / mentioned"),
    ]
    for i, (a, b, c) in enumerate(table_rows):
        ry = th_y + 0.45 + i * 0.42
        add_text(s, Inches(7.4), Inches(ry), Inches(1.9), Inches(0.4),
                 a, size=14, color=TEXT_DARK, anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(9.3), Inches(ry), Inches(1.5), Inches(0.4),
                 b, size=14, color=TEXT_GRAY, anchor=MSO_ANCHOR.MIDDLE,
                 align=PP_ALIGN.CENTER)
        add_text(s, Inches(10.8), Inches(ry), Inches(1.6), Inches(0.4),
                 c, size=14, bold=True, color=ACCENT_GREEN,
                 anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

    # ── 하단: LIMITATION 박스 (전체 폭)
    add_card(s, x=0.7, y=5.6, w=12.0, h=1.65)
    add_kicker(s, x=0.95, y=5.8, w=5, h=0.3,
               label="LIMITATION  ·  정직한 회고", color=ACCENT_WINE)
    add_text(s, Inches(0.95), Inches(6.2), Inches(11.5), Inches(0.95),
             ["•  Vercel Hobby 60초 / 일 1회 cron — backlog 누적 → "
              "Pro 업그레이드 또는 외부 워커 분리",
              "•  자막 음역 (곽재식→곽제식) → description 해시태그를 LLM 정답 표기로 주입",
              "•  배포 URL 공개 → Vercel Password Protection 또는 Supabase Auth"],
             size=13, color=TEXT_GRAY, line_spacing=1.4)

    prs.save(str(OUTPUT))
    print(f"saved: {OUTPUT.resolve()}")
    print(f"size: {OUTPUT.stat().st_size} bytes")
    print(f"slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
